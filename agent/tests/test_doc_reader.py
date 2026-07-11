"""Tests for the universal read_document tool (multi-format dispatch)."""

from __future__ import annotations

import json
from pathlib import Path

import pandas as pd
import pytest

from src.tools.doc_reader_tool import read_document


def _call(path: Path, pages: str = "") -> dict:
    return json.loads(read_document(str(path), pages))


@pytest.fixture(autouse=True)
def allow_tmp_documents(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    """Allow each test's temporary directory as a document import root."""
    monkeypatch.setenv("VIBE_TRADING_ALLOWED_FILE_ROOTS", str(tmp_path))


# ---------------- Plain text variants ----------------

@pytest.mark.parametrize("ext", [".txt", ".md", ".log", ".rst"])
def test_plain_text(tmp_path: Path, ext: str) -> None:
    p = tmp_path / f"note{ext}"
    p.write_text("hello world\nline 2", encoding="utf-8")
    result = _call(p)
    assert result["status"] == "ok"
    assert result["format"] == "text"
    assert result["encoding"] == "utf-8"
    assert "hello world" in result["text"]


@pytest.mark.parametrize("ext", [".json", ".yaml", ".yml", ".toml", ".ini", ".env"])
def test_config_formats(tmp_path: Path, ext: str) -> None:
    p = tmp_path / f"conf{ext}"
    p.write_text("key=value\n", encoding="utf-8")
    result = _call(p)
    assert result["status"] == "ok"
    assert result["format"] == "text"


@pytest.mark.parametrize("ext", [".py", ".js", ".ts", ".go", ".rs", ".sql", ".sh"])
def test_source_code(tmp_path: Path, ext: str) -> None:
    p = tmp_path / f"code{ext}"
    p.write_text("print(1)", encoding="utf-8")
    result = _call(p)
    assert result["status"] == "ok"
    assert result["format"] == "text"
    assert "print(1)" in result["text"]


def test_gbk_fallback(tmp_path: Path) -> None:
    """GBK-encoded Chinese text (common from A-share brokers)."""
    p = tmp_path / "gbk.txt"
    p.write_bytes("中文内容".encode("gbk"))
    result = _call(p)
    assert result["status"] == "ok"
    assert result["encoding"] == "gbk"
    assert "中文" in result["text"]


def test_unknown_extension_treated_as_text(tmp_path: Path) -> None:
    p = tmp_path / "weird.xyzfmt"
    p.write_text("abc", encoding="utf-8")
    result = _call(p)
    assert result["status"] == "ok"
    assert result["format"] == "text"
    assert result["text"] == "abc"


def test_csv_read_as_text(tmp_path: Path) -> None:
    csv_path = tmp_path / "sample.csv"
    csv_path.write_text("a,b,c\n1,2,3\n4,5,6\n", encoding="utf-8")
    result = _call(csv_path)
    assert result["status"] == "ok"
    assert result["format"] == "text"
    assert "a,b,c" in result["text"]


# ---------------- DOCX ----------------

def test_docx_paragraphs_and_tables(tmp_path: Path) -> None:
    import docx  # type: ignore

    p = tmp_path / "doc.docx"
    d = docx.Document()
    d.add_paragraph("First paragraph")
    d.add_paragraph("Second paragraph")
    table = d.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "A"
    table.rows[0].cells[1].text = "B"
    table.rows[1].cells[0].text = "1"
    table.rows[1].cells[1].text = "2"
    d.save(p)

    result = _call(p)
    assert result["status"] == "ok"
    assert result["format"] == "docx"
    assert result["paragraphs"] >= 2
    assert result["tables"] == 1
    assert "First paragraph" in result["text"]
    assert "A | B" in result["text"]


# ---------------- Excel ----------------

def test_xlsx_multiple_sheets(tmp_path: Path) -> None:
    p = tmp_path / "book.xlsx"
    with pd.ExcelWriter(p) as w:
        pd.DataFrame({"a": [1, 2, 3], "b": [4, 5, 6]}).to_excel(w, sheet_name="first", index=False)
        pd.DataFrame({"x": ["u", "v"]}).to_excel(w, sheet_name="second", index=False)

    result = _call(p)
    assert result["status"] == "ok"
    assert result["format"] == "excel"
    names = [s["name"] for s in result["sheets"]]
    assert names == ["first", "second"]
    assert "first" in result["text"]


# ---------------- PPTX ----------------

def test_pptx_slides(tmp_path: Path) -> None:
    from pptx import Presentation  # type: ignore

    p = tmp_path / "deck.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Alpha Title"
    prs.save(p)

    result = _call(p)
    assert result["status"] == "ok"
    assert result["format"] == "pptx"
    assert result["slides"] == 1
    assert "Alpha Title" in result["text"]


# ---------------- Error handling ----------------

def test_missing_file(tmp_path: Path) -> None:
    result = _call(tmp_path / "missing.pdf")
    assert result["status"] == "error"
    assert "not found" in result["error"].lower()


def test_directory_rejected(tmp_path: Path) -> None:
    result = _call(tmp_path)
    assert result["status"] == "error"
    assert "not a file" in result["error"].lower()


def test_envelope_shape(tmp_path: Path) -> None:
    """Every ok response must include the unified envelope keys."""
    p = tmp_path / "x.txt"
    p.write_text("hello", encoding="utf-8")
    result = _call(p)
    for key in ("status", "file", "format", "char_count", "truncated", "text"):
        assert key in result


def test_truncation(tmp_path: Path) -> None:
    p = tmp_path / "big.txt"
    p.write_text("x" * 20000, encoding="utf-8")
    result = _call(p)
    assert result["truncated"] is True
    assert result["char_count"] == 20000
    assert len(result["text"]) < 20000
